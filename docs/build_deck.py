"""Generate the demo-day presentation deck (docs/WTP_Digital_Twin.pptx).

Reproducible: edit this file and re-run  `python docs/build_deck.py`.
Theme: "Ocean" — deep blue / teal / midnight with a mint accent (fits water +
industrial control). Header font Trebuchet MS, body Calibri.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---- palette ---------------------------------------------------------------
MIDNIGHT = RGBColor(0x21, 0x29, 0x5C)
DEEPBLUE = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
MINT     = RGBColor(0x02, 0xC3, 0x9A)
LIGHT    = RGBColor(0xF4, 0xF7, 0xFA)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x21, 0x29, 0x5C)
MUTED    = RGBColor(0x5A, 0x6B, 0x7B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CORAL    = RGBColor(0xD8, 0x5A, 0x30)
AMBER    = RGBColor(0xBA, 0x75, 0x17)

HEAD = "Trebuchet MS"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


# ---- helpers ---------------------------------------------------------------
def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.0):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (txt,size,color,bold,font,italic)."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, *rest) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
            if rest and rest[0]:
                r.font.italic = True
    return tb


def title(s, txt, kicker=None):
    box(s, 0.55, 0.62, 0.12, 0.95, fill=MINT, shape=MSO_SHAPE.RECTANGLE)  # left motif bar
    if kicker:
        text(s, 0.85, 0.55, 11.5, 0.32, [[(kicker.upper(), 12, TEAL, True, HEAD)]])
        text(s, 0.83, 0.86, 11.5, 0.8, [[(txt, 30, INK, True, HEAD)]])
    else:
        text(s, 0.83, 0.7, 11.5, 0.9, [[(txt, 32, INK, True, HEAD)]])


def footer(s, n):
    text(s, 0.83, 7.0, 8, 0.3, [[("TUMA206 · Water Treatment Plant Digital Twin", 9, MUTED, False, BODY)]])
    text(s, 12.0, 7.0, 0.9, 0.3, [[(str(n), 9, MUTED, True, BODY)]], align=PP_ALIGN.RIGHT)


def node(s, l, t, w, h, head, sub, fill, txtcolor=WHITE, subcolor=None):
    box(s, l, t, w, h, fill=fill)
    subcolor = subcolor or txtcolor
    text(s, l + 0.12, t, w - 0.24, h, [
        [(head, 13, txtcolor, True, HEAD)],
        [(sub, 10.5, subcolor, False, BODY)],
    ], anchor=MSO_ANCHOR.MIDDLE, space_after=1)


def arrow(s, l, t, w, h, color=TEAL, shape=MSO_SHAPE.DOWN_ARROW):
    a = box(s, l, t, w, h, fill=color, shape=shape)
    return a


# ============================================================ SLIDE 1 — title
s = slide(MIDNIGHT)
box(s, 0, 0, SW, 0.18, fill=MINT, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 7.32, SW, 0.18, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.9, 2.15, 11.5, 0.4, [[("MODERN DEVELOPMENTS IN INDUSTRY · TUMA206", 14, MINT, True, HEAD)]])
text(s, 0.85, 2.6, 11.6, 1.6, [[("Water Treatment Plant", 50, WHITE, True, HEAD)],
                               [("Digital Twin", 50, WHITE, True, HEAD)]], space_after=2)
text(s, 0.9, 4.7, 11.5, 0.9, [[("A full industrial stack — field → PLC → MQTT → historian → "
                                "dashboard — with an agentic AI operator assistant.", 16,
                                RGBColor(0xCA, 0xDC, 0xFC), False, BODY)]])
text(s, 0.9, 6.4, 11.5, 0.5, [[("Team: [names · matriculation · cohort]   ·   Demo day 25–26 Jun 2026",
                                12, RGBColor(0x9A, 0xB0, 0xCC), False, BODY)]])

# ============================================================ SLIDE 2 — brief
s = slide(); title(s, "What we built", "The capstone brief")
intro = ("A working twin of a water-treatment tank that behaves like a real plant under "
         "stress — not just on the happy path. Every layer of the industrial stack is "
         "present, and an AI assistant helps the operator respond to faults.")
text(s, 0.85, 1.95, 11.6, 0.8, [[(intro, 15, MUTED, False, BODY)]], line_spacing=1.1)
cards = [
    ("Process twin", "Tank, pumps, valve, level & pH sensors, dosing loop — physics integrated each scan."),
    ("Deterministic PLC", "on/off filling, PID dosing and a state machine on a fixed scan cycle."),
    ("End-to-end data path", "Tag namespace → MQTT → InfluxDB historian → Grafana dashboard."),
    ("Agentic AI assistant", "Detects faults and recommends operator actions in plain language."),
]
cw, gap, x0, y0 = 5.55, 0.4, 0.85, 3.0
for i, (h, d) in enumerate(cards):
    cx = x0 + (i % 2) * (cw + gap); cy = y0 + (i // 2) * 1.85
    box(s, cx, cy, cw, 1.6, fill=CARD, line=RGBColor(0xDD, 0xE5, 0xEC))
    box(s, cx, cy, 0.1, 1.6, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx + 0.3, cy + 0.22, cw - 0.5, 1.2,
         [[(h, 16, INK, True, HEAD)], [(d, 12.5, MUTED, False, BODY)]], space_after=4, line_spacing=1.05)
footer(s, 2)

# ============================================================ SLIDE 3 — arch
s = slide(); title(s, "Architecture", "Field to AI in one stack")
# fault box
node(s, 4.7, 1.55, 4.0, 0.55, "Fault injection — operator panel", "one fault per layer, at runtime",
     RGBColor(0x88, 0x87, 0x80))
arrow(s, 6.6, 2.12, 0.2, 0.24, color=RGBColor(0x88, 0x87, 0x80))
# loop row
node(s, 1.6, 2.4, 4.2, 0.78, "Digital twin (process)", "tank · pumps · valve · pH", CORAL)
node(s, 7.6, 2.4, 4.2, 0.78, "Soft PLC", "on/off · PID · state machine", AMBER)
arrow(s, 5.8, 2.66, 1.8, 0.28, color=TEAL, shape=MSO_SHAPE.LEFT_RIGHT_ARROW)
text(s, 5.8, 2.42, 1.8, 0.22, [[("sensors ⇄ commands", 9, MUTED, False, BODY)]], align=PP_ALIGN.CENTER)
arrow(s, 6.6, 3.24, 0.2, 0.32, color=TEAL)
text(s, 6.85, 3.26, 2.0, 0.3, [[("publish tags", 9.5, MUTED, False, BODY)]])
# mqtt
node(s, 4.7, 3.66, 4.0, 0.66, "MQTT broker", "Mosquitto · pub/sub", DEEPBLUE)
arrow(s, 3.0, 4.42, 0.2, 0.3, color=TEAL); arrow(s, 10.1, 4.42, 0.2, 0.3, color=TEAL)
# branches
node(s, 1.0, 4.8, 4.0, 0.66, "Historian → InfluxDB", "time-series store", TEAL)
node(s, 8.3, 4.8, 4.0, 0.66, "AI assistant → advice", "detect · recommend · explain", MIDNIGHT)
arrow(s, 2.85, 5.56, 0.2, 0.3, color=TEAL)
node(s, 1.0, 5.94, 4.0, 0.6, "Grafana dashboard", "live trends + alarms", TEAL)
footer(s, 3)

# ============================================================ SLIDE 4 — stack
s = slide(); title(s, "Stack & justification", "LO6 — defend every choice")
rows = [
    ("Twin + PLC", "Python, single scan loop", "Transparent control logic; easy to defend in the viva."),
    ("Messaging", "MQTT (Mosquitto)", "Lightweight pub/sub; the de-facto IIoT transport; decouples field from IT."),
    ("Historian", "InfluxDB 2.7", "Purpose-built time-series DB; retention & downsampling; native Grafana link."),
    ("Dashboard", "Grafana", "No frontend code; provisioned trends and alarm thresholds."),
    ("AI assistant", "Claude (Anthropic SDK)", "Strong agentic reasoning; safe, explained advice; offline fallback."),
    ("Orchestration", "docker-compose", "One command brings up the whole stack — reproducible demo."),
]
y = 2.0; rh = 0.78
# header row
box(s, 0.85, y, 2.4, rh, fill=DEEPBLUE); box(s, 3.3, y, 3.0, rh, fill=DEEPBLUE); box(s, 6.35, y, 6.1, rh, fill=DEEPBLUE)
for cx, w, lab in [(0.85, 2.4, "Layer"), (3.3, 3.0, "Choice"), (6.35, 6.1, "Why")]:
    text(s, cx + 0.2, y, w - 0.3, rh, [[(lab, 13, WHITE, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
y += rh
for i, (a, b, c) in enumerate(rows):
    fill = CARD if i % 2 == 0 else RGBColor(0xEA, 0xF1, 0xF6)
    box(s, 0.85, y, 11.6, rh, fill=fill, line=RGBColor(0xDD, 0xE5, 0xEC))
    text(s, 1.05, y, 2.2, rh, [[(a, 12.5, INK, True, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.5, y, 2.8, rh, [[(b, 12.5, TEAL, True, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.55, y, 5.7, rh, [[(c, 11.5, MUTED, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh
footer(s, 4)

# ============================================================ SLIDE 5 — control
s = slide(); title(s, "The twin & its control", "LO2 — three control styles")
text(s, 0.85, 1.9, 11.6, 0.6, [[("One PLC scan cycle reads sensors, runs the logic, writes "
     "actuators — at a fixed period, like real hardware.", 14, MUTED, False, BODY)]], line_spacing=1.1)
items = [
    ("On/off", "Inlet pump fills the tank with hysteresis between low/high level setpoints."),
    ("PID", "Dosing pump regulates pH to setpoint with anti-windup integral clamping."),
    ("State machine", "FILLING → TREATING → DISCHARGING phases sequence the batch."),
]
y = 2.9
for h, d in items:
    box(s, 0.85, y, 0.55, 0.55, fill=MINT, shape=MSO_SHAPE.OVAL)
    text(s, 0.85, y, 0.55, 0.55, [[(h[0], 18, MIDNIGHT, True, HEAD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.65, y - 0.02, 10.6, 0.7, [[(h + "  ", 16, INK, True, HEAD), (d, 13, MUTED, False, BODY)]],
         line_spacing=1.05)
    y += 1.0
box(s, 0.85, 6.2, 11.6, 0.8, fill=RGBColor(0xEA, 0xF1, 0xF6), line=RGBColor(0xDD, 0xE5, 0xEC))
text(s, 1.1, 6.2, 11.1, 0.8, [[("Twin physics: ", 13, TEAL, True, BODY),
     ("level from a mass balance (inflow − outflow), pH a first-order lag toward the dosing "
      "equilibrium — stable and explainable.", 13, MUTED, False, BODY)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, 5)

# ============================================================ SLIDE 6 — datapath
s = slide(); title(s, "Tags & the data path", "LO3 — field to dashboard")
text(s, 0.85, 1.9, 11.6, 0.55, [[("Every value is a named tag with a unit and type — the single "
     "source of truth all services share.", 14, MUTED, False, BODY)]], line_spacing=1.1)
# tag examples
box(s, 0.85, 2.7, 6.0, 4.0, fill=CARD, line=RGBColor(0xDD, 0xE5, 0xEC))
text(s, 1.1, 2.85, 5.6, 0.4, [[("Example tags", 14, INK, True, HEAD)]])
tags = [("LIT101_level", "tank level · m"), ("AIT201_ph", "tank pH"),
        ("FIT101_inlet_flow", "influent flow · m³/s"), ("P101_inlet_cmd / _run_fb", "command + feedback"),
        ("P201_dose_cmd", "PID dosing output 0–1"), ("PLC_state", "state-machine phase"),
        ("LIT101_predicted", "mass-balance estimate")]
ty = 3.35
for name, desc in tags:
    text(s, 1.1, ty, 5.6, 0.34, [[(name, 12, TEAL, True, "Consolas"), ("   " + desc, 11, MUTED, False, BODY)]])
    ty += 0.46
# flow column
fx = 7.4
steps = [("Field device", CORAL), ("Soft PLC", AMBER), ("MQTT broker", DEEPBLUE),
         ("InfluxDB historian", TEAL), ("Grafana dashboard", TEAL)]
sy = 2.7
for i, (lab, col) in enumerate(steps):
    node(s, fx, sy, 4.6, 0.62, lab, "", col)
    sy += 0.72
    if i < len(steps) - 1:
        arrow(s, fx + 2.1, sy - 0.12, 0.18, 0.22, color=col)
footer(s, 6)

# ============================================================ SLIDE 7 — faults
s = slide(); title(s, "Fault catalog", "LO4 — one fault per Purdue layer")
rows = [
    ("Sensor", "Level transmitter stuck", "Diverges from mass-balance prediction", CORAL),
    ("Equipment", "Inlet pump runs, no flow", "Command + feedback ON while flow ≈ 0", AMBER),
    ("Process", "pH excursion", "Setpoint deviation beyond alarm band", TEAL),
    ("Infrastructure", "MQTT broker / data loss", "Telemetry staleness watchdog", DEEPBLUE),
]
y = 2.1; rh = 1.0
box(s, 0.85, y, 2.6, 0.6, fill=MIDNIGHT); box(s, 3.5, y, 3.9, 0.6, fill=MIDNIGHT); box(s, 7.45, y, 5.0, 0.6, fill=MIDNIGHT)
for cx, w, lab in [(0.85, 2.6, "Layer"), (3.5, 3.9, "Fault"), (7.45, 5.0, "Detection signal")]:
    text(s, cx + 0.2, y, w - 0.3, 0.6, [[(lab, 13, WHITE, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
y += 0.6
for layer, fault, det, col in rows:
    box(s, 0.85, y, 11.6, rh, fill=CARD, line=RGBColor(0xDD, 0xE5, 0xEC))
    box(s, 0.85, y, 0.12, rh, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, 1.15, y, 2.4, rh, [[(layer, 13, INK, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.65, y, 3.7, rh, [[(fault, 13, col, True, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 7.6, y, 4.7, rh, [[(det, 12, MUTED, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y += rh
footer(s, 7)

# ============================================================ SLIDE 8 — detection
s = slide(); title(s, "Detection approach", "Fast, deterministic, explainable")
left = [
    ("Analytical redundancy", "A stuck sensor diverges from an independent estimate (mass balance). A dead pump leaves both flat — that's how we tell them apart."),
    ("Command vs effect", "A running motor contactor doesn't prove flow. Comparing command to measured flow catches mechanical failure."),
]
right = [
    ("Setpoint deviation + debounce", "pH alarms only after the deviation persists, so transient noise doesn't trip false alarms."),
    ("Staleness on time, not value", "When all tags go quiet at once it's a data outage — not a process alarm. Operators must not chase a phantom."),
]
for col, items in [(0.85, left), (6.85, right)]:
    cy = 2.1
    for h, d in items:
        box(s, col, cy, 5.6, 1.95, fill=CARD, line=RGBColor(0xDD, 0xE5, 0xEC))
        box(s, col, cy, 5.6, 0.1, fill=MINT, shape=MSO_SHAPE.RECTANGLE)
        text(s, col + 0.3, cy + 0.25, 5.0, 1.6,
             [[(h, 16, INK, True, HEAD)], [(d, 13, MUTED, False, BODY)]], space_after=6, line_spacing=1.1)
        cy += 2.15
footer(s, 8)

# ============================================================ SLIDE 9 — AI agent
s = slide(); title(s, "The AI operator assistant", "LO5 — recommend & explain, safely")
steps = [("1", "Detect", "Rule-based checks raise a fault from sensor cross-checks and trends — independent of the LLM."),
         ("2", "Recommend", "Claude receives the fault + a recent tag snapshot and returns diagnosis, actions and reasoning as JSON."),
         ("3", "Explain & advise", "Plain-language guidance grounded in the data — the agent advises, a human decides.")]
x = 0.85
for num, h, d in steps:
    box(s, x, 2.1, 3.75, 2.5, fill=CARD, line=RGBColor(0xDD, 0xE5, 0xEC))
    box(s, x + 0.25, 2.35, 0.7, 0.7, fill=MIDNIGHT, shape=MSO_SHAPE.OVAL)
    text(s, x + 0.25, 2.35, 0.7, 0.7, [[(num, 22, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.3, 3.25, 3.2, 1.3, [[(h, 17, TEAL, True, HEAD)], [(d, 12.5, MUTED, False, BODY)]],
         space_after=6, line_spacing=1.1)
    x += 4.0
box(s, 0.85, 5.0, 11.6, 1.5, fill=MIDNIGHT)
text(s, 1.2, 5.0, 11.0, 1.5, [
    [("Safety & resilience  ", 15, MINT, True, HEAD),
     ("— the agent never actuates; it must ground claims in the data it was shown; and if the API "
      "is unavailable it falls back to a built-in playbook so the live demo always produces advice.",
      14, RGBColor(0xCA, 0xDC, 0xFC), False, BODY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s, 9)

# ============================================================ SLIDE 10 — demo
s = slide(); title(s, "Demo day flow", "Detect + alert + advise ≤ 60 s")
text(s, 0.85, 1.95, 11.6, 0.5, [[("For each of the four faults, injected live from the operator panel:",
     14, MUTED, False, BODY)]])
steps = [("Inject", "operator panel publishes the fault to the control topic"),
         ("Twin reacts", "physics + PLC respond; the fault propagates through the stack"),
         ("Detect & alarm", "rule engine raises the fault; Grafana shows the trend break"),
         ("AI recommends", "assistant posts diagnosis + actions, all within 60 seconds")]
y = 2.7
for i, (h, d) in enumerate(steps):
    box(s, 0.85, y, 11.6, 0.92, fill=CARD if i % 2 == 0 else RGBColor(0xEA, 0xF1, 0xF6),
        line=RGBColor(0xDD, 0xE5, 0xEC))
    box(s, 0.85, y, 0.92, 0.92, fill=TEAL)
    text(s, 0.85, y, 0.92, 0.92, [[(str(i + 1), 24, WHITE, True, HEAD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.1, y, 10.0, 0.92, [[(h + "  —  ", 16, INK, True, HEAD), (d, 13.5, MUTED, False, BODY)]],
         anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02
footer(s, 10)

# ============================================================ SLIDE 11 — decisions
s = slide(); title(s, "Design decisions & trade-offs", "Expect viva questions here")
pts = [
    ("Detect deterministically, then ask the LLM", "Guarantees the ≤60 s requirement even if the model is slow or offline; the LLM adds the reasoning, not the trigger."),
    ("MQTT over OPC-UA", "Simpler and lighter; we accept losing OPC-UA's richer typed metadata."),
    ("First-order physics, not CFD", "Enough fidelity to make control loops and faults realistic, while staying explainable."),
    ("InfluxDB over plain SQL", "Time-series native: retention, downsampling and Grafana integration out of the box."),
    ("Agent advises, never actuates", "Human-in-the-loop is the safe pattern for an operator assistant."),
]
y = 2.0
for h, d in pts:
    box(s, 0.85, y, 0.18, 0.86, fill=MINT, shape=MSO_SHAPE.RECTANGLE)
    text(s, 1.25, y, 11.0, 0.9, [[(h, 15, INK, True, HEAD)], [(d, 12.5, MUTED, False, BODY)]],
         space_after=2, line_spacing=1.03)
    y += 1.0
footer(s, 11)

# ============================================================ SLIDE 12 — closing
s = slide(MIDNIGHT)
box(s, 0, 0, SW, 0.18, fill=MINT, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 7.32, SW, 0.18, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.9, 2.7, 11.5, 1.2, [[("Thank you", 46, WHITE, True, HEAD)]])
text(s, 0.92, 3.9, 11.5, 0.9, [[("Questions on any part of the system — we can walk through the code.",
     16, RGBColor(0xCA, 0xDC, 0xFC), False, BODY)]])
text(s, 0.92, 5.6, 11.5, 0.5, [[("Repo: github.com/[team]/wtp-digital-twin   ·   "
     "Live demo: docker compose up", 13, MINT, False, "Consolas")]])

# ---- save ------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "WTP_Digital_Twin.pptx")
prs.save(out)
print("wrote", out, "—", len(prs.slides._sldIdLst), "slides")
